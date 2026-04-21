# ==========================================================
# 📂 NARYAN AI — FILE LOADER
#    Extracts text from PDF, DOCX, PPTX, TXT files.
# ==========================================================

import os
import logging
from typing import List, Dict, Tuple

logger = logging.getLogger("NARYAN_AI.file_loader")


def load_file(path: str) -> str:
    """Extract all text from a file (PDF / DOCX / PPTX / TXT)."""
    ext = os.path.splitext(path)[1].lower()
    try:
        if ext == ".pdf":
            return _load_pdf(path)
        elif ext == ".docx":
            return _load_docx(path)
        elif ext == ".pptx":
            return _load_pptx(path)
        elif ext in (".txt", ".md"):
            with open(path, encoding="utf-8", errors="ignore") as f:
                return f.read()
        else:
            logger.warning(f"Unsupported file type: {ext}")
            return ""
    except Exception as e:
        logger.error(f"Failed to load {path}: {e}")
        return ""


def chunk_text(text: str, chunk_size: int = 300, overlap: int = 50) -> List[str]:
    """Split text into overlapping word-level chunks."""
    words  = text.split()
    chunks = []
    step   = chunk_size - overlap
    for i in range(0, len(words), step):
        chunk = " ".join(words[i : i + chunk_size])
        if chunk.strip():
            chunks.append(chunk)
    return chunks


def load_folder(folder: str, subject: str = "", doc_type: str = "notes") -> List[Tuple[str, dict]]:
    """
    Walk a folder, load every supported file, chunk it, and return
    a list of (chunk_text, metadata) tuples ready for the vector store.
    """
    results: List[Tuple[str, dict]] = []
    supported = {".pdf", ".docx", ".pptx", ".txt", ".md"}

    for root, _, files in os.walk(folder):
        for fname in files:
            if os.path.splitext(fname)[1].lower() not in supported:
                continue
            fpath = os.path.join(root, fname)

            # Infer subject from sub-folder name if not provided
            rel   = os.path.relpath(root, folder)
            subj  = subject or (rel if rel != "." else os.path.basename(folder))

            text   = load_file(fpath)
            chunks = chunk_text(text)
            for chunk in chunks:
                results.append((chunk, {
                    "source":  fname,
                    "subject": subj,
                    "type":    doc_type,
                    "path":    fpath,
                }))
            logger.info(f"Loaded {fname} ({len(chunks)} chunks) [{subj}]")

    return results


# ── Private helpers ────────────────────────────────────────

def _load_pdf(path: str) -> str:
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            return "\n".join(p.extract_text() or "" for p in pdf.pages)
    except ImportError:
        pass
    try:
        import PyPDF2
        with open(path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            return "\n".join(
                page.extract_text() or "" for page in reader.pages
            )
    except Exception as e:
        raise RuntimeError(f"PDF load failed: {e}")


def _load_docx(path: str) -> str:
    try:
        import docx2txt
        return docx2txt.process(path)
    except ImportError:
        pass
    try:
        from docx import Document
        doc = Document(path)
        return "\n".join(p.text for p in doc.paragraphs)
    except Exception as e:
        raise RuntimeError(f"DOCX load failed: {e}")


def _load_pptx(path: str) -> str:
    try:
        from pptx import Presentation
        prs  = Presentation(path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        raise RuntimeError(f"PPTX load failed: {e}")
